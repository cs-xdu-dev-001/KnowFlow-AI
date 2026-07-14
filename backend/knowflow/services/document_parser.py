from __future__ import annotations

import csv
import html
import json
import re
from io import BytesIO, StringIO
from pathlib import Path
from typing import Any

from fastapi import HTTPException

from ..config import (
    ALLOWED_UPLOAD_SUFFIXES,
    CHUNK_OVERLAP,
    CHUNK_SIZE,
    IMAGE_MIME_TYPES,
    IMAGE_SUFFIXES,
    MAX_UPLOAD_FILE_SIZE,
)


_SAFE_FILENAME_RE = re.compile(r"[^A-Za-z0-9._\-\u4e00-\u9fff]+")


def sanitize_upload_filename(filename: str | None) -> str:
    raw = Path(filename or "document.txt").name.replace("\x00", "").strip()
    cleaned = _SAFE_FILENAME_RE.sub("_", raw).strip(" ._")
    if not cleaned:
        cleaned = "document.txt"
    if len(cleaned) > 180:
        suffix = Path(cleaned).suffix
        stem = Path(cleaned).stem[: max(1, 180 - len(suffix))]
        cleaned = f"{stem}{suffix}"
    return cleaned


def validate_upload_file(filename: str, data: bytes) -> None:
    suffix = Path(filename).suffix.lower()
    if suffix not in ALLOWED_UPLOAD_SUFFIXES:
        allowed = ", ".join(sorted(ALLOWED_UPLOAD_SUFFIXES))
        extension = suffix or "no extension"
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {extension}. Supported types: {allowed}")
    if not data:
        raise HTTPException(status_code=400, detail="Uploaded file is empty.")
    if len(data) > MAX_UPLOAD_FILE_SIZE:
        limit_mb = MAX_UPLOAD_FILE_SIZE / 1024 / 1024
        raise HTTPException(status_code=413, detail=f"File is too large. The current limit is {limit_mb:.0f} MB.")


async def read_upload_file_with_limit(upload_file) -> bytes:
    chunks: list[bytes] = []
    total = 0
    while True:
        chunk = await upload_file.read(1024 * 1024)
        if not chunk:
            break
        total += len(chunk)
        if total > MAX_UPLOAD_FILE_SIZE:
            limit_mb = MAX_UPLOAD_FILE_SIZE / 1024 / 1024
            raise HTTPException(status_code=413, detail=f"File is too large. The current limit is {limit_mb:.0f} MB.")
        chunks.append(chunk)
    return b"".join(chunks)


def split_text(text_value: str) -> list[str]:
    cleaned = re.sub(r"\r\n?", "\n", text_value).strip()
    if not cleaned:
        return []
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            separators=["\n\n", "\n", "。", "；", "，", ".", "!", "?", ",", ";", " "],
            length_function=len,
        )
        return [item.strip() for item in splitter.split_text(cleaned) if item.strip()]
    except Exception:
        chunks: list[str] = []
        start = 0
        while start < len(cleaned):
            end = min(start + CHUNK_SIZE, len(cleaned))
            piece = cleaned[start:end].strip()
            if piece:
                chunks.append(piece)
            if end == len(cleaned):
                break
            start = max(end - CHUNK_OVERLAP, start + 1)
        return chunks


def decode_bytes(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "gbk", "big5"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="ignore")


def flatten_json(value: Any, prefix: str = "") -> list[str]:
    lines: list[str] = []
    if isinstance(value, dict):
        for key, item in value.items():
            next_prefix = f"{prefix}.{key}" if prefix else str(key)
            lines.extend(flatten_json(item, next_prefix))
    elif isinstance(value, list):
        for index, item in enumerate(value):
            next_prefix = f"{prefix}[{index}]" if prefix else f"[{index}]"
            lines.extend(flatten_json(item, next_prefix))
    else:
        lines.append(f"{prefix}: {value}" if prefix else str(value))
    return lines


def parse_table_text(text_value: str, delimiter: str = ",") -> str:
    reader = csv.reader(StringIO(text_value), delimiter=delimiter)
    lines = []
    for row in reader:
        cleaned = [cell.strip() for cell in row if cell and cell.strip()]
        if cleaned:
            lines.append(" | ".join(cleaned))
    return "\n".join(lines)


def html_to_text(text_value: str) -> str:
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(text_value, "html.parser")
        for tag in soup(["script", "style", "noscript"]):
            tag.decompose()
        return soup.get_text("\n")
    except Exception:
        text_value = re.sub(r"(?is)<(script|style).*?>.*?</\1>", "\n", text_value)
        text_value = re.sub(r"(?s)<[^>]+>", "\n", text_value)
        return html.unescape(text_value)


def rtf_to_text(text_value: str) -> str:
    text_value = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text_value)
    text_value = re.sub(r"\\[a-zA-Z]+\d* ?", " ", text_value)
    text_value = text_value.replace("{", " ").replace("}", " ")
    return re.sub(r"\s+", " ", text_value)


def extract_text_from_upload(filename: str, data: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    if suffix in IMAGE_SUFFIXES:
        mime_type = IMAGE_MIME_TYPES.get(suffix, "image/*")
        return (
            f"[Image attachment]\nFilename: {filename}\nType: {mime_type}\nSize: {len(data)} bytes\n"
            "Note: KnowFlow received this image and kept a preview. Text-only models cannot read image pixels directly; "
            "connect OCR or a vision model if you need screenshot text recognition."
        )
    raw_text_suffixes = {".txt", ".md", ".markdown", ".log", ".yaml", ".yml", ".xml"}
    if suffix in raw_text_suffixes:
        return decode_bytes(data)
    if suffix == ".json":
        text_value = decode_bytes(data)
        try:
            return "\n".join(flatten_json(json.loads(text_value)))
        except Exception:
            return text_value
    if suffix == ".csv":
        return parse_table_text(decode_bytes(data), ",")
    if suffix == ".tsv":
        return parse_table_text(decode_bytes(data), "\t")
    if suffix in {".html", ".htm"}:
        return html_to_text(decode_bytes(data))
    if suffix == ".rtf":
        return rtf_to_text(decode_bytes(data))
    if suffix == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(data))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    if suffix == ".docx":
        from docx import Document

        document = Document(BytesIO(data))
        lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    lines.append(" | ".join(cells))
        return "\n".join(lines)
    if suffix in {".xlsx", ".xlsm"}:
        from openpyxl import load_workbook

        workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
        lines: list[str] = []
        for sheet in workbook.worksheets:
            lines.append(f"# Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if cells:
                    lines.append(" | ".join(cells))
        return "\n".join(lines)
    if suffix == ".pptx":
        from pptx import Presentation

        presentation = Presentation(BytesIO(data))
        lines: list[str] = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines.append(f"# Slide {index}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines.append(shape.text.strip())
        return "\n".join(lines)
    raise HTTPException(status_code=400, detail="Supported document types include txt, md, pdf, docx, xlsx, pptx, html, json, csv, tsv, rtf, yaml, xml, and log.")